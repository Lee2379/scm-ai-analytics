from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "SCM_AI_Agent_Architecture_1page.pptx"

FONT_NAME = "Noto Sans"
RED = RGBColor(230, 0, 18)
BLACK = RGBColor(17, 17, 17)
GRAY = RGBColor(90, 90, 90)
LIGHT = RGBColor(247, 247, 247)
WHITE = RGBColor(255, 255, 255)


def set_font(run, size=12, bold=False, color=BLACK) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    r_fonts = r_pr.find(qn("a:rFonts"))
    if r_fonts is None:
        r_fonts = OxmlElement("a:rFonts")
        r_pr.insert(0, r_fonts)
    r_fonts.set(qn("a:latin"), FONT_NAME)
    r_fonts.set(qn("a:ea"), FONT_NAME)
    r_fonts.set(qn("a:cs"), FONT_NAME)


def add_text(slide, x, y, w, h, text, size=12, bold=False, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if p.runs:
            set_font(p.runs[0], size=size, bold=bold, color=color)
    return box


def add_box(slide, x, y, w, h, title, subtitle, body="", fill=WHITE, accent=RED):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.6)

    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    add_text(slide, x + 0.16, y + 0.22, w - 0.32, 0.34, title, size=13, bold=True, color=BLACK)
    add_text(slide, x + 0.16, y + 0.58, w - 0.32, 0.25, subtitle, size=9.5, bold=True, color=RED)
    if body:
        add_text(slide, x + 0.16, y + 0.95, w - 0.32, h - 1.05, body, size=9.5, color=GRAY)
    return shape


def connect(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = BLACK
    line.line.width = Pt(1.2)
    return line


def build_ppt() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RED
    top_bar.line.fill.background()

    add_text(slide, 0.58, 0.42, 7.8, 0.42, "階層型AI Agentアーキテクチャ（Hierarchical AI Agent Architecture）", size=24, bold=True)
    add_text(
        slide,
        0.6,
        0.95,
        9.8,
        0.35,
        "SCM意思決定を、需要予測・在庫ポリシー・補充推薦・店舗間移動・自然言語説明に分解",
        size=11,
        color=GRAY,
    )

    # Root agent
    add_box(
        slide,
        4.05,
        1.35,
        5.15,
        1.05,
        "SCM Manager Chat Agent",
        "統合判断・自然言語説明（Integrated reasoning & explanation）",
        "",
        fill=LIGHT,
        accent=RED,
    )

    # Middle agents
    agents = [
        (
            0.55,
            "需要予測Agent",
            "(Demand Forecast Agent)",
            "SKU・店舗別の28日需要予測\n販売履歴・曜日・季節性を反映",
        ),
        (
            3.8,
            "在庫ポリシーAgent",
            "(Inventory Policy Agent)",
            "ROP・安全在庫を計算\nService Level / Lead Timeを考慮",
        ),
        (
            7.05,
            "補充推薦Agent",
            "(Replenishment Agent)",
            "現在在庫 < ROP を検出\n推奨発注数量を算出",
        ),
        (
            10.3,
            "店舗間移動Agent",
            "(Store Transfer Agent)",
            "過剰在庫から欠品リスク店舗へ\n在庫移動を推薦",
        ),
    ]
    for x, title, subtitle, body in agents:
        add_box(slide, x, 3.05, 2.72, 1.75, title, subtitle, body)
        connect(slide, 6.62, 2.4, x + 1.36, 3.05)

    # Data layer
    add_box(
        slide,
        0.72,
        5.6,
        3.0,
        0.82,
        "入力データ",
        "(Input data)",
        "Sales / Inventory / Product / Store / Supply",
        fill=LIGHT,
        accent=BLACK,
    )
    add_box(
        slide,
        4.22,
        5.6,
        3.0,
        0.82,
        "SCM計算ロジック",
        "(SCM logic)",
        "Forecast / ROP / Safety Stock / Risk Score",
        fill=LIGHT,
        accent=BLACK,
    )
    add_box(
        slide,
        7.72,
        5.6,
        4.15,
        0.82,
        "意思決定アウトプット",
        "(Decision output)",
        "Priority SKU / Order Qty / Transfer / Explanation",
        fill=LIGHT,
        accent=BLACK,
    )

    connect(slide, 2.22, 5.6, 2.0, 4.8)
    connect(slide, 5.72, 5.6, 5.16, 4.8)
    connect(slide, 9.8, 5.6, 8.4, 4.8)

    add_text(
        slide,
        0.58,
        6.92,
        11.5,
        0.28,
        "Business value: SCM業務をAgentごとに分解し、予測から実行アクションまで接続した意思決定支援システム。",
        size=9.5,
        color=GRAY,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Agent architecture PPT written to {OUT}")


if __name__ == "__main__":
    build_ppt()
