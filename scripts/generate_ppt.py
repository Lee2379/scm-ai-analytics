from __future__ import annotations

import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))


OUT = ROOT / "outputs" / "AI_SCM_Data_Analysis_Project.pptx"
FALLBACK_OUT = ROOT / "outputs" / "AI_SCM_Data_Analysis_Project_NotoSans.pptx"
DATA = ROOT / "data"

PINE = RGBColor(35, 72, 63)
INK = RGBColor(25, 37, 44)
MUTED = RGBColor(96, 112, 119)
SAGE = RGBColor(127, 155, 131)
CLAY = RGBColor(191, 133, 105)
ICE = RGBColor(238, 246, 245)
PAPER = RGBColor(251, 250, 246)
FONT_NAME = "Noto Sans"


def apply_noto_sans(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.name = FONT_NAME
                for run in paragraph.runs:
                    run.font.name = FONT_NAME
                    r_pr = run._r.get_or_add_rPr()
                    r_fonts = r_pr.find(qn("a:rFonts"))
                    if r_fonts is None:
                        r_fonts = OxmlElement("a:rFonts")
                        r_pr.insert(0, r_fonts)
                    r_fonts.set(qn("a:latin"), FONT_NAME)
                    r_fonts.set(qn("a:ea"), FONT_NAME)
                    r_fonts.set(qn("a:cs"), FONT_NAME)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.0), Inches(11.8), Inches(0.45))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = MUTED


def add_footer(slide, page: int) -> None:
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = ICE
    line.line.fill.background()
    footer = slide.shapes.add_textbox(Inches(0.58), Inches(7.12), Inches(11.8), Inches(0.25))
    p = footer.text_frame.paragraphs[0]
    p.text = f"AI SCM Data Analysis Project | Development demo | {page}"
    p.font.size = Pt(8)
    p.font.color.rgb = MUTED


def add_bullets(slide, items, x, y, w, h, font_size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = INK
        p.space_after = Pt(9)
        p.level = 0


def add_card(slide, x, y, w, h, title, body, accent=PINE):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = ICE
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.18), Inches(w - 0.35), Inches(0.32))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.bold = True
    tp.font.size = Pt(15)
    tp.font.color.rgb = PINE
    body_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.58), Inches(w - 0.35), Inches(h - 0.7))
    bp = body_box.text_frame.paragraphs[0]
    bp.text = body
    bp.font.size = Pt(12)
    bp.font.color.rgb = MUTED


def add_metric(slide, x, y, title, value):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.85), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = ICE
    t = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.15), Inches(2.5), Inches(0.25))
    tp = t.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(10)
    tp.font.bold = True
    tp.font.color.rgb = MUTED
    v = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.42), Inches(2.5), Inches(0.44))
    vp = v.text_frame.paragraphs[0]
    vp.text = str(value)
    vp.font.size = Pt(22)
    vp.font.bold = True
    vp.font.color.rgb = PINE


def load_metrics():
    policy = pd.read_csv(DATA / "inventory_policy.csv")
    recs = pd.read_csv(DATA / "recommendations.csv")
    transfers = pd.read_csv(DATA / "transfer_recommendations.csv")
    return {
        "pairs": len(policy),
        "stockout": int((policy["stock_status"] == "Stockout Risk").sum()),
        "overstock": int((policy["stock_status"] == "Overstock").sum()),
        "order_units": int(recs["recommended_order_qty"].sum()),
        "transfers": len(transfers),
    }


def build_deck() -> None:
    metrics = load_metrics()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    title = slide.shapes.add_textbox(Inches(0.65), Inches(1.25), Inches(8.7), Inches(1.4))
    p = title.text_frame.paragraphs[0]
    p.text = "AI SCM Data Analysis Project"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = INK
    sub = slide.shapes.add_textbox(Inches(0.68), Inches(2.55), Inches(8.2), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Demand forecasting, ROP, safety stock, replenishment, and store-transfer analytics for global fashion retail SCM."
    sp.font.size = Pt(17)
    sp.font.color.rgb = MUTED
    add_card(slide, 9.45, 1.15, 3.1, 3.6, "Portfolio Target", "Global Japanese retail companies such as Fast Retailing, MUJI, AEON, Rakuten, and Amazon Japan.", SAGE)
    add_footer(slide, 1)

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, "Business Problem", "Retail SCM problems are expensive because forecast errors become operational decisions.")
    add_card(slide, 0.65, 1.55, 3.8, 3.8, "Stockout", "Lost sales, lower customer satisfaction, and missed demand signals.", CLAY)
    add_card(slide, 4.75, 1.55, 3.8, 3.8, "Overstock", "Excess inventory, markdown pressure, and poor cash conversion.", SAGE)
    add_card(slide, 8.85, 1.55, 3.8, 3.8, "Logistics Inefficiency", "Poor replenishment timing and unnecessary inter-store movement.", PINE)
    add_footer(slide, 2)

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, "Solution Architecture", "An Agentic workflow that links prediction to action.")
    steps = [
        "1. Collect sales, product, store, calendar, weather, inventory, and supply assumptions",
        "2. Forecast 28-day SKU-store demand",
        "3. Calculate ROP and safety stock per SKU-store pair",
        "4. Recommend replenishment quantities and transfer opportunities",
        "5. Explain decisions through an SCM Manager Agent",
    ]
    add_bullets(slide, steps, 0.8, 1.5, 6.2, 4.5, 16)
    add_card(slide, 7.55, 1.35, 4.8, 1.1, "Demand Forecast Agent", "Predicts future SKU-store demand.", SAGE)
    add_card(slide, 7.55, 2.7, 4.8, 1.1, "Inventory Policy Agent", "Computes ROP, safety stock, and risk.", PINE)
    add_card(slide, 7.55, 4.05, 4.8, 1.1, "Replenishment Agent", "Recommends order and transfer actions.", CLAY)
    add_footer(slide, 3)

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, "Data and API Strategy", "Public data plus clearly documented SCM simulation is realistic for a portfolio.")
    add_card(slide, 0.75, 1.35, 3.8, 1.4, "Kaggle API", "M5 Forecasting and H&M fashion data can be downloaded when credentials are configured.", SAGE)
    add_card(slide, 4.75, 1.35, 3.8, 1.4, "Open-Meteo API", "Optional weather features such as temperature and precipitation for Japanese cities.", PINE)
    add_card(slide, 8.75, 1.35, 3.8, 1.4, "Simulated SCM Layer", "Inventory, lead time, service level, supplier region, and logistics assumptions.", CLAY)
    add_bullets(
        slide,
        [
            "The demo runs without private data or real company data.",
            "SCM assumptions are visible and can be stress-tested.",
            "API keys are optional and never hard-coded.",
        ],
        1.0,
        3.45,
        10.6,
        2.0,
        16,
    )
    add_footer(slide, 4)

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, "Inventory Policy Logic", "The system calculates ROP per product and store, not only by forecast horizon.")
    add_card(slide, 0.9, 1.55, 5.4, 1.45, "Reorder Point", "ROP = average daily demand x lead time + safety stock", PINE)
    add_card(slide, 0.9, 3.3, 5.4, 1.45, "Safety Stock", "Safety stock = demand standard deviation x Z-value x sqrt(lead time)", SAGE)
    add_card(slide, 6.95, 1.55, 5.1, 3.2, "Decision Rule", "If current inventory < ROP, the Agent recommends replenishment. Order quantity is calculated from target stock minus current inventory.", CLAY)
    add_footer(slide, 5)

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, "Demo Results", "Simulation-based results generated by the local portfolio system.")
    add_metric(slide, 0.75, 1.45, "SKU-store pairs", metrics["pairs"])
    add_metric(slide, 3.95, 1.45, "Stockout risks", metrics["stockout"])
    add_metric(slide, 7.15, 1.45, "Overstock cases", metrics["overstock"])
    add_metric(slide, 10.35, 1.45, "Order units", metrics["order_units"])
    add_card(slide, 1.0, 3.25, 5.3, 1.7, "Operational Output", "Prioritized reorder list with stock, ROP, safety stock, forecast demand, and business reason.", PINE)
    add_card(slide, 6.95, 3.25, 5.3, 1.7, "Transfer Output", f"{metrics['transfers']} candidate store-transfer recommendations generated from overstock to risk stores.", SAGE)
    add_footer(slide, 6)

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, "Why This Fits Global Japanese Retail", "The project connects data science with SCM business execution.")
    add_bullets(
        slide,
        [
            "Shows demand forecasting, inventory policy, and business action in one workflow",
            "Targets SCM pain points: stockout, overstock, replenishment timing, and store transfer",
            "Uses AI Agent as decision support, not as a generic chatbot",
            "Can be localized for Japan with holidays, weather, city/store assumptions, and Japanese responses",
            "Demonstrates communication between analytics, engineering, and business stakeholders",
        ],
        0.9,
        1.4,
        11.4,
        4.4,
        17,
    )
    add_footer(slide, 7)

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, "Next Steps", "How to make the portfolio even stronger before submission.")
    add_card(slide, 0.8, 1.45, 3.7, 3.5, "1. Real Dataset Extension", "Connect Kaggle API and replace synthetic sales with M5/H&M data transformations.", SAGE)
    add_card(slide, 4.8, 1.45, 3.7, 3.5, "2. Logistics API", "Add Google Routes API for route time and transfer cost estimation.", PINE)
    add_card(slide, 8.8, 1.45, 3.7, 3.5, "3. Gemini Agent", "Use Gemini only from environment variables to explain recommendations in Japanese and English.", CLAY)
    add_footer(slide, 8)

    apply_noto_sans(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(OUT)
        print(f"Presentation written to {OUT}")
    except PermissionError:
        prs.save(FALLBACK_OUT)
        print(f"Presentation written to {FALLBACK_OUT}")


if __name__ == "__main__":
    build_deck()
