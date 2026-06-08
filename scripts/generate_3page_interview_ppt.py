from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "AI_SCM_Data_Analysis_Project_3pages.pptx"

FONT_NAME = "Noto Sans"
RED = RGBColor(230, 0, 18)
BLACK = RGBColor(17, 17, 17)
GRAY = RGBColor(96, 96, 96)
LIGHT = RGBColor(246, 246, 246)
WHITE = RGBColor(255, 255, 255)


def set_font(run, size=None, bold=None, color=None) -> None:
    run.font.name = FONT_NAME
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    r_fonts = r_pr.find(qn("a:rFonts"))
    if r_fonts is None:
        r_fonts = OxmlElement("a:rFonts")
        r_pr.insert(0, r_fonts)
    r_fonts.set(qn("a:latin"), FONT_NAME)
    r_fonts.set(qn("a:ea"), FONT_NAME)
    r_fonts.set(qn("a:cs"), FONT_NAME)


def add_text(slide, x, y, w, h, text, size=16, bold=False, color=BLACK, line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, raw_line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = raw_line
        p.line_spacing = line_spacing
        if p.runs:
            set_font(p.runs[0], size=size, bold=bold, color=color)
    return box


def add_header(slide, page_no: int, section: str) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    add_text(slide, 0.55, 0.35, 4.8, 0.28, section, size=10, bold=True, color=RED)
    add_text(slide, 11.9, 0.35, 0.8, 0.28, f"0{page_no}", size=10, bold=True, color=GRAY)


def add_footer(slide) -> None:
    add_text(
        slide,
        0.55,
        7.08,
        7.5,
        0.22,
        "AI SCM Data Analysis Project",
        size=8,
        color=GRAY,
    )
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(6.98), Inches(12.2), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(230, 230, 230)
    line.line.fill.background()


def add_screenshot_placeholder(slide, x, y, w, h, label) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = RED
    shape.line.width = Pt(1.2)
    add_text(slide, x + 0.25, y + 0.25, w - 0.5, 0.35, label, size=12, bold=True, color=RED)
    add_text(
        slide,
        x + 0.25,
        y + h / 2 - 0.28,
        w - 0.5,
        0.6,
        "Insert screenshot here",
        size=20,
        bold=True,
        color=GRAY,
    )


def add_red_label(slide, x, y, text) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(1.05), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    add_text(slide, x + 0.12, y + 0.07, 0.85, 0.18, text, size=8, bold=True, color=WHITE)


def build_ppt() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Page 1
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, 1, "Project Overview")
    add_red_label(slide, 0.65, 1.0, "SCM AI")
    add_text(
        slide,
        0.65,
        1.45,
        5.6,
        1.15,
        "AI SCM Data Analysis Project",
        size=27,
        bold=True,
        color=BLACK,
        line_spacing=0.9,
    )
    add_text(
        slide,
        0.68,
        2.85,
        5.35,
        0.55,
        "AI SCM\u30c7\u30fc\u30bf\u5206\u6790\u30d7\u30ed\u30b8\u30a7\u30af\u30c8",
        size=13,
        bold=True,
        color=RED,
        line_spacing=0.95,
    )
    add_text(
        slide,
        0.68,
        3.72,
        5.6,
        1.65,
        "本プロジェクトは、欠品リスク、過剰在庫、補充判断、店舗間在庫移動の課題に対して、データサイエンスとAI Agentを活用したSCM意思決定支援システムを構築したものです。\n\n需要予測、SKU・店舗別の発注点、安全在庫、推奨発注数量を一つのダッシュボードに統合し、SCM担当者が迅速に判断できる仕組みを設計しました。",
        size=12.5,
        color=BLACK,
        line_spacing=1.08,
    )
    add_screenshot_placeholder(slide, 6.75, 1.05, 5.9, 5.55, "Screenshot 1: Dashboard overview")
    add_footer(slide)

    # Page 2
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, 2, "Inventory Logic")
    add_text(
        slide,
        0.65,
        0.9,
        5.7,
        0.85,
        "SKU・店舗別の発注点（ROP）と安全在庫の設計",
        size=25,
        bold=True,
        color=BLACK,
    )
    add_text(
        slide,
        0.68,
        1.95,
        5.55,
        0.9,
        "単純な需要予測だけではなく、SKU・店舗単位で発注点（ROP）と安全在庫を計算し、実際の補充判断につなげています。",
        size=12.5,
        color=BLACK,
    )
    add_text(slide, 0.68, 3.05, 5.3, 0.35, "ROP = 日平均需要 × リードタイム + 安全在庫", size=17, bold=True, color=RED)
    add_text(slide, 0.68, 3.55, 5.65, 0.5, "安全在庫 = 需要の標準偏差 × Z値（サービスレベル）× √リードタイム", size=14, bold=True, color=BLACK)
    add_text(
        slide,
        0.68,
        4.45,
        5.6,
        1.6,
        "判断ロジック\n1. 現在在庫 < ROP のSKUを検出\n2. 欠品リスクと需要予測をもとに優先順位を設定\n3. 推奨発注数量を算出\n4. 必要に応じて店舗間在庫移動を検討",
        size=12.5,
        color=BLACK,
        line_spacing=1.08,
    )
    add_screenshot_placeholder(slide, 6.75, 1.05, 5.9, 5.55, "Screenshot 2: ROP & Safety Stock table")
    add_footer(slide)

    # Page 3
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, 3, "AI Agent Recommendation")
    add_text(
        slide,
        0.65,
        0.9,
        5.7,
        0.85,
        "AI Agentによる補充判断と説明可能な意思決定支援",
        size=25,
        bold=True,
        color=BLACK,
    )
    add_text(
        slide,
        0.68,
        1.9,
        5.6,
        1.45,
        "AI Chatbotは、ダッシュボード上のSCMデータをもとに、どのSKUを優先的に再発注すべきかを自然言語で説明します。\n\n在庫、ROP、安全在庫、28日需要予測、推奨発注数量を参照し、SCM担当者が理解しやすい形で判断根拠を提示します。",
        size=12.5,
        color=BLACK,
        line_spacing=1.08,
    )
    add_text(slide, 0.68, 3.85, 5.4, 0.48, "Example Question", size=11, bold=True, color=RED)
    add_text(
        slide,
        0.68,
        4.32,
        5.6,
        0.9,
        "どのSKUを優先的に再発注すべきかを教えてください。\nまた、その判断ロジックについても説明してください。",
        size=13,
        bold=True,
        color=BLACK,
    )
    add_text(
        slide,
        0.68,
        5.62,
        5.6,
        0.8,
        "Expected Output: 優先SKU Top 3、判断ロジック、次のアクションを要約表示",
        size=12,
        bold=True,
        color=RED,
    )
    add_screenshot_placeholder(slide, 6.75, 1.05, 5.9, 5.55, "Screenshot 3: AI Chatbot answer")
    add_footer(slide)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"3-page PPT written to {OUT}")


if __name__ == "__main__":
    build_ppt()
