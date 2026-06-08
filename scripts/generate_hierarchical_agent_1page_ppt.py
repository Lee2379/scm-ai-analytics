"""1-page presentation deck: hierarchical AI agent architecture with a white and red clean theme."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "AI_SCM_Data_Analysis_Project_Agent_Architecture_1page.pptx"

FONT_NAME = "Noto Sans"
RED = RGBColor(230, 0, 18)
BLACK = RGBColor(17, 17, 17)
GRAY = RGBColor(96, 96, 96)
LIGHT = RGBColor(246, 246, 246)
WHITE = RGBColor(255, 255, 255)
LINE_GRAY = RGBColor(230, 230, 230)


def set_font(run, size=None, bold=None, color=None) -> None:
    run.font.name = FONT_NAME
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    r_fonts = r_pr.find(qn("a:rFonts"))
    if r_fonts is None:
        r_fonts = OxmlElement("a:rFonts")
        r_pr.insert(0, r_fonts)
    for tag in ("a:latin", "a:ea", "a:cs"):
        r_fonts.set(qn(tag), FONT_NAME)


def add_text(slide, x, y, w, h, text, size=16, bold=False, color=BLACK, line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, raw_line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = raw_line
        p.line_spacing = line_spacing
        if p.runs:
            set_font(p.runs[0], size=size, bold=bold, color=color)
    return box


def add_header(slide, page_no: int, section: str) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    add_text(slide, 0.55, 0.35, 6.5, 0.28, section, size=10, bold=True, color=RED)
    add_text(slide, 11.9, 0.35, 0.8, 0.28, f"{page_no:02d}", size=10, bold=True, color=GRAY)


def add_footer(slide) -> None:
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(6.98), Inches(12.2), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE_GRAY
    line.line.fill.background()
    add_text(
        slide,
        0.55,
        7.08,
        12.0,
        0.22,
        "AI SCM Data Analysis Project",
        size=8,
        color=GRAY,
    )


def add_diagram_frame(slide, x, y, w, h, label: str) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = RED
    shape.line.width = Pt(1.2)
    add_text(slide, x + 0.22, y + 0.18, w - 0.44, 0.32, label, size=11, bold=True, color=RED)


def node(slide, x, y, w, h, title_jp: str, title_en: str, fill=WHITE):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = RED
    box.line.width = Pt(1.0)
    stripe = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.07))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RED
    stripe.line.fill.background()
    add_text(
        slide,
        x + 0.08,
        y + 0.12,
        w - 0.16,
        h - 0.14,
        f"{title_jp}\n{title_en}",
        size=8.2,
        bold=True,
        color=BLACK,
        line_spacing=0.95,
    )


def connector(slide, x1, y1, x2, y2) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = GRAY
    line.line.width = Pt(1.0)


def build_diagram(slide) -> None:
    fx, fy, fw, fh = 6.75, 1.05, 5.9, 5.55
    add_diagram_frame(slide, fx, fy, fw, fh, "Hierarchical Agent Map / 階層型Agent構成図")

    cx = fx + fw / 2
    # L1
    node(slide, cx - 1.55, fy + 0.62, 3.1, 0.62, "SCMマネージャー Agent", "Orchestrator", fill=WHITE)
    connector(slide, cx, fy + 1.24, cx, fy + 1.42)

    # L2
    node(slide, fx + 0.35, fy + 1.48, 2.35, 0.58, "Gemini 3.5 Flash", "LLM Agent (Primary)")
    node(slide, fx + 3.2, fy + 1.48, 2.35, 0.58, "ローカル・ルール Agent", "Rule Agent (Fallback)")
    connector(slide, cx, fy + 1.42, fx + 1.52, fy + 1.48)
    connector(slide, cx, fy + 1.42, fx + 4.37, fy + 1.48)
    connector(slide, cx, fy + 2.06, cx, fy + 2.22)

    # L3 specialists
    specs = [
        (fx + 0.22, "需要予測", "Forecast"),
        (fx + 1.72, "在庫ポリシー", "ROP / SS"),
        (fx + 3.22, "補充推薦", "Replenish"),
        (fx + 4.72, "店舗間移動", "Transfer"),
    ]
    for i, (sx, jp, en) in enumerate(specs):
        node(slide, sx, fy + 2.28, 1.35, 0.72, f"{jp} Agent", en)
        connector(slide, cx, fy + 2.22, sx + 0.67, fy + 2.28)

    # L4 foundation
    node(
        slide,
        fx + 0.35,
        fy + 3.35,
        fw - 0.7,
        0.55,
        "データ基盤 → scm_engine.py",
        "Sales · Inventory · Forecast · Policy · Recommendations",
        fill=LIGHT,
    )
    connector(slide, cx, fy + 3.0, cx, fy + 3.35)

    add_text(
        slide,
        fx + 0.28,
        fy + 4.15,
        fw - 0.56,
        1.2,
        "Flow: CSV → Engine → Context → Gemini or Local → NL Answer\n"
        "Decision logic: reproducible calculations with a natural-language explanation layer",
        size=8.5,
        color=GRAY,
        line_spacing=1.1,
    )


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_header(slide, 1, "Project Title & Agent Architecture")

    add_text(
        slide,
        0.65,
        0.9,
        5.75,
        0.9,
        "AI SCM Data Analysis Project",
        size=24,
        bold=True,
        color=BLACK,
        line_spacing=0.92,
    )
    add_text(
        slide,
        0.68,
        1.72,
        5.55,
        0.45,
        "AI SCM\u30c7\u30fc\u30bf\u5206\u6790\u30d7\u30ed\u30b8\u30a7\u30af\u30c8",
        size=13,
        bold=True,
        color=RED,
        line_spacing=0.95,
    )
    add_text(
        slide,
        0.68,
        2.15,
        5.55,
        1.15,
        "本システムは、SCM業務を専門Agentに分解し、需要予測から補充・店舗間移動まで一貫した意思決定支援を行います。\n"
        "チャット層はライブCSVコンテキストを参照し、数値ロジックと説明可能な回答を提供します。",
        size=12.5,
        color=BLACK,
        line_spacing=1.08,
    )
    add_text(
        slide,
        0.68,
        3.35,
        5.5,
        0.35,
        "Architecture Layers / 構成レイヤー",
        size=11,
        bold=True,
        color=RED,
    )
    add_text(
        slide,
        0.68,
        3.72,
        5.65,
        1.55,
        "L1  Orchestrator: SCM Manager Agent（質問受付・統合説明）\n"
        "L2  Reasoning: Gemini 3.5 Flash（Primary）/ Local Rule Agent（Fallback）\n"
        "L3  Specialists: Forecast · Inventory Policy · Replenishment · Transfer\n"
        "L4  Foundation: scm_engine + demo CSV → Dashboard & Agent Context",
        size=12,
        color=BLACK,
        line_spacing=1.12,
    )
    add_text(
        slide,
        0.68,
        5.55,
        5.65,
        0.75,
        "Business Value:\n"
        "Hierarchical workflow from forecast to ROP, replenishment, transfer, and explanation",
        size=11.5,
        bold=True,
        color=RED,
        line_spacing=1.05,
    )

    build_diagram(slide)
    add_footer(slide)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Written: {OUT}")


if __name__ == "__main__":
    build()
